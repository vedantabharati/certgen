import os
import sys
import shutil
import pandas as pd
import subprocess
import tempfile
import re
import argparse
from pptx import Presentation
from concurrent.futures import ProcessPoolExecutor, as_completed
import multiprocessing

# Path to LibreOffice binary on macOS
SOFFICE_PATH = "/opt/homebrew/bin/soffice"

def to_title_case(text):
    return str(text).strip().title()

def safe_filename(name):
    # Remove any characters that aren't letters, numbers, underscores, or hyphens
    return re.sub(r'[^a-zA-Z0-9_\-\.]', '', name.replace(' ', '_'))

def generate_pptx(args_tuple):
    """
    Worker function to generate a single PPTX file.
    Takes a tuple to work easily with ProcessPoolExecutor.
    """
    index, row_data, template_path, default_school, temp_dir = args_tuple
    
    name = to_title_case(row_data.get('Name', ''))
    
    # Check if school is in CSV, otherwise use default flag
    school = ''
    if 'School' in row_data and pd.notna(row_data['School']) and str(row_data['School']).strip():
        school = to_title_case(row_data['School'])
    elif default_school:
        school = default_school
    
    # Load PPTX
    prs = Presentation(template_path)
    
    # Replace Tags
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if '<<Full Name>>' in run.text:
                            run.text = run.text.replace('<<Full Name>>', name)
                        if '<<School>>' in run.text:
                            run.text = run.text.replace('<<School>>', school)

    # Save Temp File securely. Prefix with index to prevent overwriting identical names.
    safe_name = safe_filename(name)
    temp_pptx = os.path.join(temp_dir, f"{index}_{safe_name}.pptx")
    prs.save(temp_pptx)
    
    return temp_pptx

def convert_batch(args_tuple):
    """
    Worker function to convert a batch of PPTX files to PDF using a unique LibreOffice profile.
    This allows us to run multiple LibreOffice instances in parallel!
    """
    chunk, outdir = args_tuple
    
    # Create an isolated user profile so parallel LibreOffice instances do not clash
    with tempfile.TemporaryDirectory() as profile_dir:
        profile_url = f"file://{profile_dir}"
        cmd = [
            SOFFICE_PATH, 
            f'-env:UserInstallation={profile_url}',
            '--headless', 
            '--convert-to', 'pdf',
            '--outdir', outdir
        ] + chunk
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"❌ Error during conversion batch:")
            print(result.stderr)
        return len(chunk)

def main():
    parser = argparse.ArgumentParser(description="Generate Certificates from a CSV and a PPTX template")
    parser.add_argument("data_csv", help="Path to the input CSV file containing student names")
    parser.add_argument("--template", required=True, help="Path to the PPTX template file")
    parser.add_argument("--school", default="", help="Default school name to fill in <<School>> if it is missing from the CSV")
    parser.add_argument("--limit", type=int, default=None, help="Limit the number of generation processes (useful for testing)")
    parser.add_argument("--outdir", default="./Certificates", help="Output directory for generated PDFs")
    
    args = parser.parse_args()
    
    print("🚀 Booting up high-performance certificate generator...")
    
    # 0. Sanity Checks
    if not os.path.exists(SOFFICE_PATH):
        print(f"❌ Error: LibreOffice not found at {SOFFICE_PATH}")
        print("Please install LibreOffice or verify the SOFFICE_PATH.")
        return

    if not os.path.exists(args.template):
        print(f"❌ Error: Template file '{args.template}' not found.")
        return

    try:
        df = pd.read_csv(args.data_csv)
    except FileNotFoundError:
        print(f"❌ Error: Data CSV '{args.data_csv}' not found.")
        return

    if 'Name' not in df.columns:
        print(f"❌ Error: Data CSV '{args.data_csv}' must contain a 'Name' column.")
        return

    if args.limit and args.limit > 0:
        print(f"⚠️ Limiting execution to first {args.limit} records as requested.")
        df = df.head(args.limit)

    # Directories setup dynamically based on user flag
    temp_dir = os.path.join(args.outdir, "temp_pptx")
    
    if not os.path.exists(args.outdir): 
        os.makedirs(args.outdir)
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)

    # 1. Parallel PPTX Generation
    print(f"🛠️  Generating step 1: Creating PPTX files for {len(df)} records...")
    generated_files = []
    
    # Pass index to prevent filename collisions
    worker_args = [(index, row, args.template, args.school, temp_dir) for index, row in df.iterrows()]
    
    max_workers = multiprocessing.cpu_count()
    
    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(generate_pptx, w_arg): w_arg for w_arg in worker_args}
        done_count = 0
        for future in as_completed(futures):
            try:
                temp_pptx = future.result()
                generated_files.append(temp_pptx)
            except Exception as e:
                print(f"❌ Error generating presentation: {e}")
            done_count += 1
            if done_count % 50 == 0 or done_count == len(worker_args):
                print(f"✨ Generated {done_count}/{len(worker_args)} PPTX files...")

    # 2. Parallel PDF Conversion (Batching)
    print("\n🚀 Generating step 2: Starting parallel PDF conversion...")
    
    # Group into chunks. A chunk of 20 balances start-up overhead and concurrency
    chunk_size = 20 
    chunks = [generated_files[i:i + chunk_size] for i in range(0, len(generated_files), chunk_size)]
    batch_args = [(chunk, args.outdir) for chunk in chunks]
    
    # Using multiple processors to convert separate chunks concurrently
    with ProcessPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(convert_batch, b_arg) for b_arg in batch_args]
        converted_count = 0
        for future in as_completed(futures):
            try:
                count = future.result()
                if count:
                    converted_count += count
                    print(f"📄 Converted {converted_count}/{len(generated_files)} to PDF...")
            except Exception as e:
                print(f"❌ Error in conversion batch: {e}")

    # 3. Cleanup
    print("\n🧹 Cleaning up temporary `.pptx` files...")
    try:
        shutil.rmtree(temp_dir)
    except Exception as e:
        print(f"⚠️ Could not fully clean up temp files: {e}")

    print(f"✅ Done! Check the {args.outdir} folder.")

if __name__ == "__main__":
    # Prevent infinite loop error on macOS multiprocessing
    multiprocessing.freeze_support()
    main()
